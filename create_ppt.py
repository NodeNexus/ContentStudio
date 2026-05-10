from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DECK_PATH = "ContentStudio_Project_Deck.pptx"

BG_DARK = RGBColor(11, 12, 20)
BG_PANEL = RGBColor(20, 22, 40)
PRIMARY = RGBColor(124, 58, 237)
ACCENT = RGBColor(34, 211, 238)
TEXT = RGBColor(245, 246, 255)
MUTED = RGBColor(180, 182, 214)


def add_background(slide):
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()


def add_title(slide, title, subtitle=""):
    add_background(slide)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.font.name = "Inter"
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(11.8), Inches(0.9))
        stf = sub.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(17)
        sp.font.color.rgb = MUTED
        sp.font.name = "Inter"


def add_bullets(slide, items, top=2.1):
    body = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.8), Inches(4.9))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(23)
        p.font.color.rgb = TEXT
        p.font.name = "Inter"
        p.space_after = Pt(11)


def add_two_col(slide, left_title, left_items, right_title, right_items):
    panel_left = slide.shapes.add_shape(1, Inches(0.7), Inches(1.7), Inches(6.1), Inches(5.2))
    panel_left.fill.solid()
    panel_left.fill.fore_color.rgb = BG_PANEL
    panel_left.line.color.rgb = PRIMARY
    panel_right = slide.shapes.add_shape(1, Inches(6.9), Inches(1.7), Inches(5.7), Inches(5.2))
    panel_right.fill.solid()
    panel_right.fill.fore_color.rgb = BG_PANEL
    panel_right.line.color.rgb = ACCENT

    lt = slide.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.5), Inches(0.6))
    lt.text_frame.text = left_title
    lt.text_frame.paragraphs[0].font.size = Pt(21)
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.color.rgb = TEXT
    rt = slide.shapes.add_textbox(Inches(7.2), Inches(1.95), Inches(5.1), Inches(0.6))
    rt.text_frame.text = right_title
    rt.text_frame.paragraphs[0].font.size = Pt(21)
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.color.rgb = TEXT

    lb = slide.shapes.add_textbox(Inches(1.0), Inches(2.55), Inches(5.5), Inches(4.1)).text_frame
    rb = slide.shapes.add_textbox(Inches(7.2), Inches(2.55), Inches(5.1), Inches(4.1)).text_frame
    lb.word_wrap = True
    rb.word_wrap = True
    lb.clear()
    rb.clear()
    for i, item in enumerate(left_items):
        p = lb.paragraphs[0] if i == 0 else lb.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED
    for i, item in enumerate(right_items):
        p = rb.paragraphs[0] if i == 0 else rb.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED


def add_footer(slide, text="Content Studio | 2026"):
    foot = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12), Inches(0.3))
    tf = foot.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Content Studio", "Autonomous Multi-Agent AI Marketing Platform")
    add_bullets(
        s,
        [
            "Production-grade SaaS experience with cinematic UI/UX",
            "Built with Flask + CrewAI + SQLite + real-time pipeline visualization",
            "Transforms one brief into a complete multi-channel content package",
        ],
        top=2.7,
    )
    add_footer(s, "Prepared by NodeNexus Team")

    # Slide 2: Problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "The Problem", "Why content teams need agentic automation")
    add_bullets(
        s,
        [
            "Marketing teams spend most of their time creating, not strategizing",
            "Single-model workflows cannot specialize per channel effectively",
            "Cross-platform consistency is hard under tight launch timelines",
            "Manual coordination across research, writing, social, email, and video is expensive",
        ],
    )
    add_footer(s)

    # Slide 3: Solution
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Our Solution", "An AI marketing team in one platform")
    add_bullets(
        s,
        [
            "Input: topic, audience, tone",
            "Orchestrated multi-agent execution using CrewAI",
            "Output: research brief, blog, social set, email sequence, video script",
            "Export-ready in Markdown, JSON, TXT with user history and dashboard analytics",
        ],
    )
    add_footer(s)

    # Slide 4: Agent Team
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "AI Agent Team", "Five specialist agents collaborating in pipeline")
    add_two_col(
        s,
        "Research & Strategy",
        [
            "Research Analyst: gathers insights, trends, key points",
            "Context-aware brief passed to downstream agents",
            "Improves factual grounding and direction",
        ],
        "Content Production",
        [
            "Blog Writer: long-form SEO content",
            "Social Media Specialist: channel-native posts",
            "Email Marketing Expert: conversion-first campaigns",
            "Video Scriptwriter: high-retention scripts",
        ],
    )
    add_footer(s)

    # Slide 5: Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Technical Architecture", "End-to-end full-stack design")
    add_bullets(
        s,
        [
            "Frontend: Flask templates + premium design system (glassmorphism, motion, accessibility)",
            "Backend: Flask API routes, CrewAI orchestration, Flask-Login authentication",
            "Database: SQLite for user and generation history persistence",
            "Integrations: LLM provider via Groq/OpenAI-compatible endpoint",
            "Deployment-ready structure with Render configuration and modular static assets",
        ],
    )
    add_footer(s)

    # Slide 6: Workflow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Workflow Orchestration", "How one prompt becomes campaign-ready output")
    add_bullets(
        s,
        [
            "1) User submits topic + audience + tone",
            "2) Research agent executes first and prepares strategic context",
            "3) Parallel specialized agents produce channel-specific assets",
            "4) System aggregates output, computes stats, stores history",
            "5) User previews, copies, and exports content package instantly",
        ],
    )
    add_footer(s)

    # Slide 7: UI/UX
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Product UI/UX", "Modern SaaS interface ready for investor demos")
    add_bullets(
        s,
        [
            "Design language inspired by Linear, Vercel, OpenAI, Perplexity",
            "Dark premium palette with glow accents and depth layering",
            "Responsive pages: Landing, Auth, Dashboard, Workspace, Visualizer, History, Settings",
            "Microinteractions: animated nodes, typing indicators, hover elevation, motion-safe mode",
            "Accessibility: focus-visible states, keyboard-friendly modal/tabs, reduced-motion support",
        ],
    )
    add_footer(s)

    # Slide 8: Key Features
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Key Features", "What makes Content Studio production-grade")
    add_two_col(
        s,
        "Core Product Features",
        [
            "Authenticated user workspaces",
            "Real-time pipeline visualizer",
            "Live markdown content viewer",
            "History storage and reload",
            "Multi-format export modal",
        ],
        "Operational Capabilities",
        [
            "Cache support for repeat prompts",
            "Usage and content stat summaries",
            "Mobile-first responsive layouts",
            "Scalable CSS token architecture",
            "Fast iteration with Flask template system",
        ],
    )
    add_footer(s)

    # Slide 9: Demo Flow
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Demo Flow", "Recommended 3-minute showcase sequence")
    add_bullets(
        s,
        [
            "Open landing + dashboard to explain value proposition",
            "Enter prompt in workspace and trigger generation",
            "Show live AI node states and activity log progression",
            "Review generated assets across tabs (blog/social/email/video/research)",
            "Export package and load historical run from database",
        ],
    )
    add_footer(s)

    # Slide 10: Performance & Value
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Business Value & Performance", "Why this matters to teams and stakeholders")
    add_bullets(
        s,
        [
            "2-3 minute average turnaround for complete content package",
            "10x productivity potential for content operations",
            "Consistent brand-aligned messaging across channels",
            "From ideation to deliverable in a single interface",
            "Clear upgrade path to enterprise integrations and analytics",
        ],
    )
    add_footer(s)

    # Slide 11: Security and Reliability
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Security & Reliability", "Production considerations already in place")
    add_bullets(
        s,
        [
            "User authentication and session management with Flask-Login",
            "Persistent history and audit-friendly generated timestamps",
            "Robust API error handling and user feedback toasts",
            "Source-controlled deployment config and repeatable startup instructions",
        ],
    )
    add_footer(s)

    # Slide 12: Roadmap
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Roadmap", "Next milestones for scale and monetization")
    add_bullets(
        s,
        [
            "Team workspaces with project collaboration and approvals",
            "Publishing integrations (WordPress, LinkedIn, X, email platforms)",
            "Template libraries and brand voice profiles",
            "A/B testing and performance analytics loop",
            "Enterprise plans: SSO, RBAC, and usage governance",
        ],
    )
    add_footer(s)

    # Slide 13: Ask
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "The Ask", "What we are seeking next")
    add_bullets(
        s,
        [
            "Pilot partners for real-world content workflows",
            "Product feedback on team collaboration and integrations",
            "Support for scaling from demo-ready to multi-tenant SaaS",
        ],
        top=2.5,
    )
    add_footer(s)

    # Slide 14: Closing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Thank You", "Content Studio — Your Autonomous AI Marketing Team")
    add_bullets(
        s,
        [
            "GitHub: https://github.com/NodeNexus/ContentStudio",
            "Built with Flask, CrewAI, and a premium product-first UI system",
            "Ready for Product Hunt, investor demos, and pilot onboarding",
        ],
        top=2.8,
    )
    add_footer(s, "Contact: NodeNexus Team")

    prs.save(DECK_PATH)
    print(f"Saved: {DECK_PATH}")


if __name__ == "__main__":
    build_deck()
